from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Create presentation
prs = Presentation()

# Set slide dimensions (16:9 widescreen)
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def add_slide_with_title(prs, layout_idx, title_text):
    """Add a slide with a title"""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    return slide

def format_title(title_shape):
    """Format title with RBC colors"""
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 194)  # RBC Blue
    
def add_bullet_points(text_frame, points):
    """Add bullet points to a text frame"""
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

# SLIDE 1: Title Slide
slide = add_slide_with_title(prs, 0, "RBC Digital Financial Wellness Platform")
subtitle = slide.placeholders[1]
subtitle.text = "18-Month Implementation Plan\n$10M Investment → $25M Revenue\nSteering Committee Presentation"

# SLIDE 2: Executive Summary
slide = add_slide_with_title(prs, 1, "Executive Summary")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Strategic Investment: $10M over 18 months",
    "Revenue Target: $25M incremental revenue",
    "Customer Impact: 200,000+ active users within 6 months",
    "Business Value: 20% increase in cross-sell conversion rates",
    "Timeline: 3 phases over 18 months with controlled rollout",
    "Risk Mitigation: Comprehensive contingency plans included"
])

# SLIDE 3: Vision & Strategic Value
slide = add_slide_with_title(prs, 1, "Vision & Strategic Value")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Transform customer financial wellness journey",
    "Position RBC as leader in financial health",
    "Leverage real RBC data for immediate customer value",
    "Create natural cross-selling opportunities",
    "Build trust through value-first approach",
    "Generate sustainable competitive advantage"
])

# SLIDE 4: 18-Month Roadmap Overview
slide = add_slide_with_title(prs, 1, "18-Month Implementation Roadmap")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Phase 1 (Months 1-3): Requirements & Architecture",
    "  • Stakeholder workshops and alignment",
    "  • Technical foundation and compliance framework",
    "Phase 2 (Months 4-15): Agile Development",
    "  • Real data integration and core features",
    "  • Customer co-creation and continuous testing",
    "Phase 3 (Months 16-18): Controlled Launch",
    "  • Pilot with 1,000 → 25,000 → Full rollout"
])

# SLIDE 5: Phase 1 Details - Requirements & Architecture
slide = add_slide_with_title(prs, 1, "Phase 1: Requirements & Architecture (Months 1-3)")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Month 1: Stakeholder Working Sessions",
    "  • Business requirements with all 15+ teams",
    "  • Technical architecture design",
    "Month 2: Validation & Compliance",
    "  • Customer advisory panel (25 members)",
    "  • Regulatory framework (PIPEDA, OSFI)",
    "Month 3: Foundation Setup",
    "  • API design and security implementation",
    "  • 500 employee volunteer testing program"
])

# SLIDE 6: Phase 2 - Development Sprints
slide = add_slide_with_title(prs, 1, "Phase 2: Agile Development (Months 4-15)")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Sprint Block 1 (M4-6): Core App & Real Data",
    "  • SSO integration, transaction analysis, budgeting",
    "Sprint Block 2 (M7-9): Insights & Education",
    "  • Personalized insights, financial education, analytics",
    "Sprint Block 3 (M10-12): Savings & Goals",
    "  • Gamification, automated savings, health score",
    "Sprint Block 4 (M13-15): Support & Optimization",
    "  • AI chatbot, advanced features, pre-launch testing"
])

# SLIDE 7: Phase 3 - Controlled Launch
slide = add_slide_with_title(prs, 1, "Phase 3: Controlled Launch (Months 16-18)")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Month 16: Limited Pilot",
    "  • 1,000 customers → 5,000 customers",
    "  • Success criteria: 70% weekly active users",
    "Month 17: Market Testing",
    "  • Expand to 25,000 customers",
    "  • Load testing and optimization",
    "Month 18: Full Launch",
    "  • Rollout to entire customer base",
    "  • Integration roadmap for main RBC app"
])

# SLIDE 8: Key Features & Cross-Selling
slide = add_slide_with_title(prs, 1, "Key Features & Revenue Generation")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Real-Time Data Integration",
    "  • Instant value with existing RBC account data",
    "Gamified Savings Challenges",
    "  • Natural introduction to RBC savings products",
    "Financial Education Platform",
    "  • Contextual product recommendations",
    "Goal-Based Planning",
    "  • Investment opportunity identification",
    "Financial Health Score",
    "  • Personalized improvement paths"
])

# SLIDE 9: Stakeholder Engagement
slide = add_slide_with_title(prs, 1, "Stakeholder Engagement Strategy")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Executive Sponsors: Monthly steering committee",
    "Business Units: Weekly sprint reviews",
    "Customer Advisory Panel: Bi-weekly feedback",
    "Compliance Team: Embedded in development",
    "Employee Testing: 500 volunteer program",
    "Branch Network: Change management support"
])

# SLIDE 10: Risk Management
slide = add_slide_with_title(prs, 1, "Risk Management & Mitigation")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Technical Risks",
    "  • API backup methods, cloud scaling protocols",
    "Business Risks",
    "  • Pivot criteria: <60% adoption triggers review",
    "Timeline Risks",
    "  • 2-week buffers in each sprint block",
    "Compliance Risks",
    "  • Embedded officers, continuous audits",
    "Security Risks",
    "  • Penetration testing, incident response plan"
])

# SLIDE 11: Financial Analysis
slide = add_slide_with_title(prs, 1, "Financial Analysis & ROI")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Investment Breakdown ($10M)",
    "  • 60% Personnel, 25% Technology, 10% External, 5% Buffer",
    "Revenue Projections ($25M)",
    "  • $15M from increased deposits",
    "  • $7M from product cross-sell",
    "  • $3M from cost avoidance",
    "Break-Even: Month 14",
    "ROI: 2.5x return over 18 months"
])

# SLIDE 12: Success Metrics
slide = add_slide_with_title(prs, 1, "Success Metrics & KPIs")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Customer Metrics",
    "  • 200,000+ active users",
    "  • 65% monthly active rate",
    "  • 4.5+ app store rating",
    "Business Metrics",
    "  • 20% increase in cross-sell conversion",
    "  • 15% increase in digital engagement",
    "Operational Metrics",
    "  • 99.9% uptime, <2% defect rate",
    "  • 95% customer support satisfaction"
])

# SLIDE 13: Demo & Live Application
slide = add_slide_with_title(prs, 1, "Live Demo - See It In Action")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Desktop Application Demo",
    "  • https://nagaraj1991.github.io/digital-banking-transformation_RBC-Project/",
    "",
    "Mobile Application Demo",
    "  • https://nagaraj1991.github.io/digital-banking-transformation_RBC-Project/mobile.html",
    "",
    "Key Features to Demonstrate:",
    "  • Real-time account integration simulation",
    "  • Financial wellness dashboard",
    "  • Goals and savings tracking"
])

# SLIDE 14: Next Steps
slide = add_slide_with_title(prs, 1, "Immediate Next Steps")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Upon Approval (48 hours):",
    "  • Form core project team",
    "  • Schedule stakeholder workshops",
    "  • Initiate technical architecture review",
    "Week 1 Deliverables:",
    "  • Detailed project charter",
    "  • Stakeholder engagement calendar",
    "  • Customer advisory panel recruitment",
    "Required from Steering Committee:",
    "  • Executive sponsor designation",
    "  • Budget authorization"
])

# SLIDE 15: Call to Action
slide = add_slide_with_title(prs, 1, "Investment Decision")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "The Opportunity:",
    "  • $10M investment → $25M revenue",
    "  • Market leadership in financial wellness",
    "  • Enhanced customer relationships",
    "",
    "The Ask:",
    "  • Approval for $10M investment",
    "  • Executive sponsorship commitment",
    "  • Authorization to proceed with Phase 1",
    "",
    "The Time is Now: Competitive window closing"
])

# SLIDE 16: Appendix
slide = add_slide_with_title(prs, 1, "Appendix - Additional Resources")
format_title(slide.shapes.title)
content = slide.placeholders[1].text_frame
add_bullet_points(content, [
    "Detailed Documentation Available:",
    "  • Complete 18-month project plan",
    "  • Risk register with mitigation strategies",
    "  • Technical architecture specifications",
    "  • Compliance framework details",
    "  • Customer research findings",
    "  • Competitive analysis report",
    "  • Financial models and assumptions",
    "",
    "Questions? Contact: Project Team"
])

# Save the presentation
output_path = "/home/nagarajan2219/digital-banking-transformation/presentation/RBC_Digital_Wellness_Platform.pptx"
prs.save(output_path)
print(f"PowerPoint presentation created successfully!")
print(f"Location: {output_path}")
print(f"Total slides: 16")
print("\nSlide Overview:")
print("1. Title Slide")
print("2. Executive Summary")
print("3. Vision & Strategic Value")
print("4. 18-Month Roadmap Overview")
print("5. Phase 1: Requirements & Architecture")
print("6. Phase 2: Agile Development")
print("7. Phase 3: Controlled Launch")
print("8. Key Features & Revenue Generation")
print("9. Stakeholder Engagement Strategy")
print("10. Risk Management & Mitigation")
print("11. Financial Analysis & ROI")
print("12. Success Metrics & KPIs")
print("13. Live Demo Links")
print("14. Immediate Next Steps")
print("15. Investment Decision (Call to Action)")
print("16. Appendix")